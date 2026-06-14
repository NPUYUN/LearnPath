from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path
from urllib.parse import urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from app.models.schemas import ClassroomHandoutSection, ClassroomSessionResponse, ClassroomSlide
from app.services.media_storage import media_image_path


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
INK = RGBColor(17, 24, 39)
MUTED = RGBColor(82, 95, 117)
LINE = RGBColor(222, 229, 238)
PAPER = RGBColor(250, 250, 247)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(35, 102, 247)
GREEN = RGBColor(0, 137, 102)
AMBER = RGBColor(191, 100, 12)
RED = RGBColor(214, 55, 73)
PURPLE = RGBColor(101, 75, 230)


ACCENTS: dict[str, RGBColor] = {
    "blue": BLUE,
    "teal": RGBColor(0, 128, 128),
    "amber": AMBER,
    "indigo": RGBColor(76, 81, 191),
    "green": GREEN,
    "rose": RED,
    "violet": PURPLE,
    "cyan": RGBColor(0, 132, 180),
}


def _set_run_typeface(run, font_name: str = "Microsoft YaHei") -> None:
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        element = r_pr.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            r_pr.append(element)
        element.set("typeface", font_name)


def _accent(item: ClassroomSlide | None = None) -> RGBColor:
    if not item:
        return BLUE
    return ACCENTS.get((item.accent_color or "blue").strip(), BLUE)


def _clean(text: object, limit: int | None = None) -> str:
    value = re.sub(r"\s+", " ", str(text or "").strip())
    if limit and len(value) > limit:
        return value[: limit - 1] + "…"
    return value


def _split_items(value: object, limit: int = 5) -> list[str]:
    if isinstance(value, list):
        raw = [str(x) for x in value]
    else:
        raw = re.split(r"[\n；;。]", str(value or ""))
    items: list[str] = []
    for item in raw:
        cleaned = _strip_list_marker(_clean(item))
        if len(cleaned) >= 2 and cleaned not in items:
            items.append(cleaned)
    return items[:limit]


def _strip_list_marker(value: str) -> str:
    text = value.strip()
    text = re.sub(r"^\s*[-*]\s+", "", text)
    text = re.sub(r"^\s*\d{1,2}[)）、:：]\s*", "", text)
    text = re.sub(r"^\s*\d{1,2}\.\s+", "", text)
    text = re.sub(r"^\s*\d{1,2}\s+(?=[\u4e00-\u9fffA-Za-z（(])", "", text)
    return text.strip()


def _is_weak(text: str) -> bool:
    value = _clean(text)
    weak_words = ["清爽", "科技", "封面", "图标", "视觉", "展示", "课堂页", "知识结构"]
    return len(value) < 12 or sum(1 for word in weak_words if word in value) >= 2


def _local_image_path(image_url: str) -> Path | None:
    raw = (image_url or "").strip()
    if not raw:
        return None
    parsed = urlparse(raw)
    path = parsed.path or raw
    prefix = "/api/media/images/"
    if prefix not in path:
        return None
    filename = path.split(prefix, 1)[1].split("/", 1)[0]
    if not filename or ".." in filename or "\\" in filename:
        return None
    local = media_image_path(filename)
    return local if local.is_file() else None


def _fill(shape, color: RGBColor, line: RGBColor | None = None) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line or color


def _textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size: int,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font_name: str = "Microsoft YaHei",
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text or ""
    _set_run_typeface(run, font_name)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _rich_text(slide, left, top, width, height, paragraphs: list[str], size: int = 15, color: RGBColor = INK):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, text in enumerate(paragraphs):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(8)
        run = paragraph.add_run()
        run.text = text
        _set_run_typeface(run)
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return shape


def _section_label(slide, left, top, label: str, accent: RGBColor) -> None:
    box = slide.shapes.add_shape(1, left, top, Inches(1.2), Inches(0.28))
    _fill(box, accent)
    _textbox(slide, left + Inches(0.08), top + Inches(0.035), Inches(1.0), Inches(0.18), label, 8, WHITE, True, PP_ALIGN.CENTER)


def _add_header(slide, item: ClassroomSlide, index: int, total: int, accent: RGBColor) -> None:
    _textbox(slide, Inches(0.55), Inches(0.28), Inches(2.4), Inches(0.24), _clean(item.kicker or f"{index:02d} / LESSON", 24), 9, accent, True)
    _textbox(slide, Inches(11.55), Inches(0.28), Inches(1.2), Inches(0.24), f"{index:02d}/{total:02d}", 9, MUTED, True, PP_ALIGN.RIGHT)
    rule = slide.shapes.add_shape(1, Inches(0.55), Inches(0.62), Inches(12.25), Inches(0.01))
    _fill(rule, LINE)


def _card(slide, left, top, width, height, fill: RGBColor = WHITE, line: RGBColor = LINE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    _fill(shape, fill, line)
    return shape


def _callout(slide, left, top, width, height, label: str, text: str, accent: RGBColor) -> None:
    _card(slide, left, top, width, height, RGBColor(245, 248, 255), RGBColor(205, 216, 255))
    _textbox(slide, left + Inches(0.18), top + Inches(0.15), width - Inches(0.36), Inches(0.22), label, 9, accent, True)
    _textbox(slide, left + Inches(0.18), top + Inches(0.45), width - Inches(0.36), height - Inches(0.58), text, 13, INK, True)


def _table(slide, left, top, width, height, columns: list[str], rows: list[list[str]], accent: RGBColor) -> bool:
    columns = [_clean(x, 20) for x in columns if _clean(x)][:4]
    rows = [[_clean(cell, 46) for cell in row[: len(columns)]] for row in rows if isinstance(row, list)]
    rows = [row for row in rows if any(row)][:6]
    if len(columns) < 2 or not rows:
        return False
    graphic = slide.shapes.add_table(len(rows) + 1, len(columns), left, top, width, height)
    table = graphic.table
    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        cell.text_frame.margin_left = Inches(0.06)
        cell.text_frame.margin_right = Inches(0.06)
        _set_cell_text(cell, col, 9, WHITE, True, PP_ALIGN.CENTER)
    for r, row in enumerate(rows, start=1):
        for c in range(len(columns)):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(247, 249, 252)
            cell.text_frame.margin_left = Inches(0.07)
            cell.text_frame.margin_right = Inches(0.07)
            _set_cell_text(cell, row[c] if c < len(row) else "", 8, INK, False)
    return True


def _set_cell_text(cell, text: str, size: int, color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    cell.text_frame.clear()
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    _set_run_typeface(run)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _latex_value(block: dict) -> str:
    for key in ("latex", "formula", "expression"):
        value = _clean(block.get(key))
        if value:
            return value
    return ""


def _format_latex(value: str) -> str:
    value = _clean(value, 180)
    if not value:
        return ""
    value = value.strip("$")

    def frac(match: re.Match[str]) -> str:
        numerator = match.group(1)
        denominator = match.group(2)
        return f"({numerator})/({denominator})"

    value = re.sub(r"\\frac\{([^{}]+)\}\{([^{}]+)\}", frac, value)
    replacements = {
        r"\to": "→",
        r"\lim": "lim",
        r"\cdot": "·",
        r"\times": "×",
        r"\Delta": "Δ",
        r"\infty": "∞",
        r"\left": "",
        r"\right": "",
    }
    for old, new in replacements.items():
        value = value.replace(old, new)
    value = value.replace("{", "").replace("}", "")
    value = re.sub(r"\s+", " ", value).strip()
    return value


def _mathtext_formula(value: str) -> str:
    formula = _clean(value, 220).strip()
    while formula.startswith("$") and formula.endswith("$") and len(formula) >= 2:
        formula = formula[1:-1].strip()
    formula = formula.replace(r"\dfrac", r"\frac").replace(r"\tfrac", r"\frac")
    formula = formula.replace(r"\displaystyle", "")
    formula = formula.replace(r"\left", "").replace(r"\right", "")
    return formula.strip()


def _format_latex(value: str) -> str:
    value = _mathtext_formula(value)
    if not value:
        return ""

    def frac(match: re.Match[str]) -> str:
        numerator = match.group(1)
        denominator = match.group(2)
        return f"({numerator})/({denominator})"

    value = re.sub(r"\\frac\{([^{}]+)\}\{([^{}]+)\}", frac, value)
    replacements = {
        r"\to": "\u2192",
        r"\lim": "lim",
        r"\cdot": "\u00b7",
        r"\times": "\u00d7",
        r"\Delta": "\u0394",
        r"\infty": "\u221e",
        r"\left": "",
        r"\right": "",
        r"\displaystyle": "",
    }
    for old, new in replacements.items():
        value = value.replace(old, new)
    value = value.replace("{", "").replace("}", "")
    value = re.sub(r"\s+", " ", value).strip()
    return value


def _render_formula_png(value: str) -> BytesIO | None:
    formula = _mathtext_formula(value)
    if not formula:
        return None
    try:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig = plt.figure(figsize=(7.0, 1.0), dpi=220)
        ax = fig.add_axes([0, 0, 1, 1])
        ax.axis("off")
        ax.text(0.5, 0.5, f"${formula}$", ha="center", va="center", fontsize=26, color="#111827")
        buffer = BytesIO()
        fig.savefig(buffer, format="png", transparent=True, bbox_inches="tight", pad_inches=0.08)
        plt.close(fig)
        buffer.seek(0)
        return buffer
    except Exception:
        return None


def _formula_card(
    slide,
    left,
    top,
    width,
    height,
    block: dict,
    item: ClassroomSlide,
    accent: RGBColor,
    *,
    include_steps: bool = True,
) -> bool:
    formula = _latex_value(block)
    if not formula:
        return False
    _card(slide, left, top, width, height, RGBColor(247, 249, 252), RGBColor(205, 216, 230))
    _textbox(
        slide,
        left + Inches(0.22),
        top + Inches(0.18),
        width - Inches(0.44),
        Inches(0.28),
        _clean(block.get("title") or "公式", 34),
        10,
        accent,
        True,
    )
    band = slide.shapes.add_shape(1, left + Inches(0.22), top + Inches(0.58), width - Inches(0.44), Inches(0.72))
    _fill(band, WHITE, RGBColor(218, 226, 238))
    formula_image = _render_formula_png(formula)
    if formula_image:
        slide.shapes.add_picture(
            formula_image,
            left + Inches(0.48),
            top + Inches(0.67),
            width=width - Inches(0.96),
            height=Inches(0.48),
        )
    else:
        _textbox(
            slide,
            left + Inches(0.38),
            top + Inches(0.76),
            width - Inches(0.76),
            Inches(0.28),
            _format_latex(formula),
            16,
            INK,
            True,
            PP_ALIGN.CENTER,
            "Cambria Math",
        )
    explanation = _clean(block.get("explanation") or block.get("answer") or item.teacher_note, 160)
    if explanation:
        _textbox(slide, left + Inches(0.26), top + Inches(1.48), width - Inches(0.52), Inches(0.64), explanation, 11, MUTED)
    steps = _split_items(block.get("steps") or block.get("items") or item.board, 4)
    if include_steps and steps and height > Inches(2.9):
        _steps(slide, left + Inches(0.22), top + Inches(2.28), width - Inches(0.44), height - Inches(2.48), steps, accent, columns=1)
    return True


def _steps(slide, left, top, width, height, steps: list[str], accent: RGBColor, columns: int = 1) -> bool:
    steps = [_clean(x, 54) for x in steps if _clean(x)][:6]
    if not steps:
        return False
    gap = Inches(0.12)
    col_w = (width - gap * (columns - 1)) / columns
    row_h = (height - gap * ((len(steps) + columns - 1) // columns - 1)) / max(1, (len(steps) + columns - 1) // columns)
    for i, step in enumerate(steps):
        col = i % columns
        row = i // columns
        x = left + (col_w + gap) * col
        y = top + (row_h + gap) * row
        _card(slide, x, y, col_w, row_h, RGBColor(248, 250, 252), LINE)
        badge = slide.shapes.add_shape(1, x + Inches(0.12), y + Inches(0.12), Inches(0.34), Inches(0.34))
        _fill(badge, accent)
        _textbox(slide, x + Inches(0.12), y + Inches(0.17), Inches(0.34), Inches(0.12), str(i + 1), 8, WHITE, True, PP_ALIGN.CENTER)
        _textbox(slide, x + Inches(0.58), y + Inches(0.10), col_w - Inches(0.72), row_h - Inches(0.12), step, 10, INK, True)
    return True


def _visual_block_type(block: dict) -> str:
    return str(block.get("type") or "process").lower().strip()


def _slide_blocks(item: ClassroomSlide) -> list[dict]:
    return [x for x in (item.visual_blocks or []) if isinstance(x, dict)]


def _primary_block(item: ClassroomSlide, preferred: set[str] | None = None) -> dict | None:
    blocks = _slide_blocks(item)
    if preferred:
        for block in blocks:
            if _visual_block_type(block) in preferred:
                return block
    return blocks[0] if blocks else None


def _board_points(item: ClassroomSlide, limit: int = 5) -> list[str]:
    points = _split_items(item.board, limit)
    if len(points) >= 3:
        return points
    fallback = _split_items(item.body, limit)
    return (points + [x for x in fallback if x not in points])[:limit]


def _proof_rows_from_points(points: list[str]) -> list[list[str]]:
    labels = ["核心含义", "判断方式", "使用场景", "常见误区", "自查问题"]
    rows: list[list[str]] = []
    for i, point in enumerate(points[:5]):
        rows.append([labels[i] if i < len(labels) else f"要点 {i + 1}", point])
    return rows


def _render_block(slide, item: ClassroomSlide, left, top, width, height, accent: RGBColor, dense: bool = False) -> bool:
    block = _primary_block(item)
    points = _board_points(item, 5)
    if not block:
        return _steps(slide, left, top, width, height, points, accent, columns=2 if dense else 1)

    btype = _visual_block_type(block)
    if btype in {"table", "compare"}:
        if _table(slide, left, top, width, height, block.get("columns") or [], block.get("rows") or [], accent):
            return True
    if btype == "formula" or _latex_value(block):
        if _formula_card(slide, left, top, width, height, block, item, accent):
            return True
    if btype in {"example", "exercise", "formula"}:
        question = _clean(block.get("question") or block.get("expression") or item.title, 110)
        steps = _split_items(block.get("steps") or block.get("items") or points, 4)
        answer = _clean(block.get("answer") or block.get("explanation") or item.teacher_note, 120)
        _card(slide, left, top, width, height, WHITE, LINE)
        _textbox(slide, left + Inches(0.22), top + Inches(0.2), width - Inches(0.44), Inches(0.28), _clean(block.get("title") or "例题拆解", 28), 11, accent, True)
        _textbox(slide, left + Inches(0.22), top + Inches(0.58), width - Inches(0.44), Inches(0.72), question, 15, INK, True)
        _steps(slide, left + Inches(0.22), top + Inches(1.45), width - Inches(0.44), Inches(1.8), steps, accent, columns=1)
        _callout(slide, left + Inches(0.22), top + height - Inches(0.82), width - Inches(0.44), Inches(0.56), "答案检查", answer, accent)
        return True

    steps = _split_items(block.get("steps") or block.get("items") or points, 6)
    if _steps(slide, left, top, width, height, steps, accent, columns=2 if dense else 1):
        return True
    return _table(slide, left, top, width, height, ["检查项", "内容"], _proof_rows_from_points(points), accent)


def _add_cover(prs: Presentation, session: ClassroomSessionResponse, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK
    accent = BLUE
    _textbox(slide, Inches(0.75), Inches(0.58), Inches(2.4), Inches(0.28), "LEARNPATH AI CLASSROOM", 9, RGBColor(164, 185, 255), True)
    _textbox(slide, Inches(0.75), Inches(1.45), Inches(8.8), Inches(1.35), _clean(session.title, 48), 39, WHITE, True)
    _textbox(slide, Inches(0.78), Inches(3.0), Inches(7.9), Inches(0.86), _clean(session.objective, 120), 18, RGBColor(216, 224, 238))

    _card(slide, Inches(9.35), Inches(1.15), Inches(3.15), Inches(4.85), RGBColor(31, 41, 55), RGBColor(57, 72, 96))
    _textbox(slide, Inches(9.65), Inches(1.48), Inches(2.4), Inches(0.3), "自学路线", 14, WHITE, True)
    roadmap = ["问题：先知道为什么要学", "概念：建立直觉和正式表达", "例题：完整拆解一遍", "练习：换场景自己判断", "复盘：带着任务离开"]
    _steps(slide, Inches(9.65), Inches(1.95), Inches(2.55), Inches(3.2), roadmap, accent, columns=1)
    _textbox(slide, Inches(9.65), Inches(5.42), Inches(2.6), Inches(0.3), f"{total} 页课件 · 预计 {session.estimated_minutes or 20} 分钟", 10, RGBColor(187, 199, 218), True)

    sources = "、".join([x.title for x in (session.source_resources or [])[:3] if x.title]) or "根据当前学习目标和课堂画像生成"
    _callout(slide, Inches(0.75), Inches(5.72), Inches(7.4), Inches(0.7), "参考资料", _clean(sources, 80), accent)


def _add_concept_slide(prs: Presentation, item: ClassroomSlide, index: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    accent = _accent(item)
    _add_header(slide, item, index, total, accent)
    _textbox(slide, Inches(0.65), Inches(0.92), Inches(7.0), Inches(0.8), _clean(item.title, 54), 27, INK, True)
    body = _clean(item.body, 190)
    if _is_weak(body):
        body = f"这一页要把“{_clean(item.title, 32)}”讲成一个能使用的判断工具：先抓住定义，再看适用条件，最后用例子检验。"
    _rich_text(slide, Inches(0.68), Inches(1.92), Inches(5.95), Inches(1.25), [body], 15, MUTED)

    points = _board_points(item, 5)
    _section_label(slide, Inches(0.68), Inches(3.45), "必须掌握", accent)
    _steps(slide, Inches(0.68), Inches(3.85), Inches(5.8), Inches(2.42), points, accent, columns=1)

    _section_label(slide, Inches(7.05), Inches(1.08), "知识对象", accent)
    formula_block = _primary_block(item, {"formula"})
    block = _primary_block(item, {"table", "compare"})
    if formula_block and _formula_card(slide, Inches(7.05), Inches(1.52), Inches(5.55), Inches(3.1), formula_block, item, accent):
        pass
    elif block and _table(slide, Inches(7.05), Inches(1.52), Inches(5.55), Inches(3.1), block.get("columns") or [], block.get("rows") or [], accent):
        pass
    else:
        _table(slide, Inches(7.05), Inches(1.52), Inches(5.55), Inches(3.1), ["检查项", "你要能说清"], _proof_rows_from_points(points), accent)
    _callout(slide, Inches(7.05), Inches(4.92), Inches(5.55), Inches(0.95), "自学提示", _clean(item.teacher_note or "如果看不懂，先不要背术语，把每个概念改写成自己的话。", 100), accent)


def _add_example_slide(prs: Presentation, item: ClassroomSlide, index: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    accent = _accent(item)
    _add_header(slide, item, index, total, accent)
    _textbox(slide, Inches(0.65), Inches(0.92), Inches(11.6), Inches(0.65), _clean(item.title, 60), 26, INK, True)
    _textbox(slide, Inches(0.68), Inches(1.65), Inches(11.2), Inches(0.45), _clean(item.body, 145), 14, MUTED)
    _section_label(slide, Inches(0.68), Inches(2.35), "题目与解法", accent)
    _render_block(slide, item, Inches(0.68), Inches(2.78), Inches(7.05), Inches(3.6), accent, dense=False)
    _section_label(slide, Inches(8.05), Inches(2.35), "为什么这么做", accent)
    points = _board_points(item, 4)
    _steps(slide, Inches(8.05), Inches(2.78), Inches(4.55), Inches(2.38), points, accent, columns=1)
    _callout(slide, Inches(8.05), Inches(5.45), Inches(4.55), Inches(0.72), "检查标准", "你不是只看答案，而是能解释每一步的依据。", accent)


def _add_table_slide(prs: Presentation, item: ClassroomSlide, index: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    accent = _accent(item)
    _add_header(slide, item, index, total, accent)
    _textbox(slide, Inches(0.65), Inches(0.92), Inches(7.8), Inches(0.7), _clean(item.title, 58), 27, INK, True)
    _textbox(slide, Inches(0.68), Inches(1.7), Inches(7.2), Inches(0.6), _clean(item.body, 150), 14, MUTED)
    block = _primary_block(item, {"table", "compare"})
    if block and _table(slide, Inches(0.75), Inches(2.55), Inches(11.85), Inches(3.35), block.get("columns") or [], block.get("rows") or [], accent):
        pass
    else:
        _table(slide, Inches(0.75), Inches(2.55), Inches(11.85), Inches(3.35), ["错误理解", "为什么错", "正确抓法"], [
            ["只背定义", "不会判断什么时候能用", "先写适用条件"],
            ["跳过例题", "看不到概念如何落地", "至少跑通一个最小例子"],
            ["只看答案", "不知道错在第几步", "把解题步骤和依据对应起来"],
        ], accent)
    _callout(slide, Inches(0.75), Inches(6.12), Inches(11.85), Inches(0.52), "这一页怎么用", "先遮住右列自己判断，再对照正确抓法；如果说不出理由，回到上一页例题。", accent)


def _add_process_slide(prs: Presentation, item: ClassroomSlide, index: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    accent = _accent(item)
    _add_header(slide, item, index, total, accent)
    _textbox(slide, Inches(0.65), Inches(0.92), Inches(7.8), Inches(0.72), _clean(item.title, 58), 27, INK, True)
    _textbox(slide, Inches(0.68), Inches(1.75), Inches(5.9), Inches(1.0), _clean(item.body, 165), 15, MUTED)
    _section_label(slide, Inches(0.68), Inches(3.1), "学习流程", accent)
    _render_block(slide, item, Inches(0.68), Inches(3.52), Inches(11.9), Inches(2.58), accent, dense=True)
    _callout(slide, Inches(7.1), Inches(1.55), Inches(5.45), Inches(0.82), "暂停点", "每走完一步都问自己：我现在知道了什么？下一步为什么必要？", accent)


def _add_quiz_slide(prs: Presentation, item: ClassroomSlide, index: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 251, 255)
    accent = _accent(item)
    _add_header(slide, item, index, total, accent)
    _textbox(slide, Inches(0.65), Inches(0.92), Inches(8.5), Inches(0.7), _clean(item.title, 58), 28, INK, True)
    _textbox(slide, Inches(0.68), Inches(1.78), Inches(5.9), Inches(0.72), _clean(item.body, 150), 15, MUTED)
    _render_block(slide, item, Inches(0.75), Inches(2.85), Inches(6.5), Inches(3.3), accent, dense=False)
    _section_label(slide, Inches(7.85), Inches(2.85), "做完再看", accent)
    _steps(slide, Inches(7.85), Inches(3.28), Inches(4.65), Inches(1.8), ["先独立写出一句解释", "再标出已知条件和目标", "最后说出答案依据"], accent)
    _callout(slide, Inches(7.85), Inches(5.38), Inches(4.65), Inches(0.76), "过关标准", "能换一个新场景仍然说出同一套判断逻辑。", accent)


def _add_summary_slide(prs: Presentation, session: ClassroomSessionResponse, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK
    accent = GREEN
    _textbox(slide, Inches(0.75), Inches(0.55), Inches(2.0), Inches(0.28), "REVIEW", 9, RGBColor(156, 230, 198), True)
    _textbox(slide, Inches(0.75), Inches(1.0), Inches(7.8), Inches(0.8), "这节课结束后，你应该能独立完成这三件事", 28, WHITE, True)
    tasks = session.homework or ["用自己的话复述核心概念", "完成一道变式练习", "记录仍然不确定的一个问题"]
    _steps(slide, Inches(0.82), Inches(2.05), Inches(5.7), Inches(3.2), tasks[:5], accent, columns=1)
    _card(slide, Inches(7.15), Inches(1.82), Inches(5.15), Inches(3.7), RGBColor(31, 41, 55), RGBColor(57, 72, 96))
    _textbox(slide, Inches(7.48), Inches(2.12), Inches(4.55), Inches(0.3), "讲义摘要", 14, WHITE, True)
    handout = session.handout[:4] if session.handout else []
    if handout:
        lines = [f"{x.heading}：{_clean(x.content, 58)}" for x in handout]
    else:
        lines = ["把概念、例题、易错点和自测题整理成一页复盘。"]
    _rich_text(slide, Inches(7.48), Inches(2.62), Inches(4.48), Inches(2.3), lines, 12, RGBColor(218, 226, 240))
    _textbox(slide, Inches(0.82), Inches(6.42), Inches(11.3), Inches(0.26), f"LearnPath AI Classroom · {total} 页自学课件", 10, RGBColor(184, 197, 214), True)


def _slide_kind(item: ClassroomSlide) -> str:
    layout = (item.layout or "").lower()
    block_types = {_visual_block_type(x) for x in _slide_blocks(item)}
    if layout in {"example"} or block_types & {"example"}:
        return "example"
    if block_types & {"formula"}:
        return "concept"
    if layout in {"mistake"} or block_types & {"table", "compare"}:
        return "table"
    if layout in {"quiz"} or block_types & {"exercise"}:
        return "quiz"
    if layout in {"timeline", "summary"} or block_types & {"process", "diagram"}:
        return "process"
    return "concept"


def _content_score(item: ClassroomSlide) -> int:
    score = 0
    score += 1 if len(_clean(item.title)) >= 6 else 0
    score += 1 if len(_clean(item.body)) >= 30 else 0
    score += 1 if len(_board_points(item, 5)) >= 3 else 0
    blocks = _slide_blocks(item)
    score += 1 if blocks else 0
    rich_types = {_visual_block_type(x) for x in blocks}
    score += 1 if rich_types & {"table", "compare", "example", "exercise", "formula", "process"} else 0
    return score


def _repair_slide(item: ClassroomSlide, index: int, session: ClassroomSessionResponse) -> ClassroomSlide:
    points = _board_points(item, 5)
    if not points:
        points = ["先理解问题", "再掌握关键概念", "最后用例题检查"]
    if _is_weak(item.body):
        item.body = f"本页围绕“{_clean(item.title or session.title, 36)}”展开。自学时先读核心解释，再看右侧结构，最后用下方检查项确认自己是否真的理解。"
    if len(points) < 3:
        points = (points + ["说出适用条件", "完成一个最小例题", "检查是否能迁移"])[:5]
    item.board = points
    if not item.visual_blocks:
        if index % 4 == 0:
            item.visual_blocks = [{
                "type": "table",
                "title": "概念对照",
                "columns": ["问题", "判断方式", "自查"],
                "rows": [
                    ["这个概念解决什么", points[0], "能否用一句话说明"],
                    ["什么时候能用", points[1], "能否说出条件"],
                    ["怎么证明会了", points[2], "能否换题迁移"],
                ],
            }]
        elif index % 4 == 1:
            item.visual_blocks = [{
                "type": "example",
                "title": "最小例题",
                "question": f"围绕“{_clean(item.title, 28)}”设计一个最小问题，并说明怎么判断。",
                "steps": points[:4],
                "answer": "答案不只看结果，要能说明使用了哪个条件和哪一步规则。",
            }]
        else:
            item.visual_blocks = [{"type": "process", "title": "学习步骤", "steps": points[:5]}]
    return item


def _prepare_slides(session: ClassroomSessionResponse) -> list[ClassroomSlide]:
    slides = list(session.slides or [])
    if not slides:
        slides = [
            ClassroomSlide(
                title=session.title or "AI 自学课堂",
                body=session.objective or "围绕当前目标建立概念、例题和练习。",
                board=["理解问题", "掌握概念", "完成练习"],
                layout="concept",
            )
        ]
    repaired: list[ClassroomSlide] = []
    for index, item in enumerate(slides[:18], start=1):
        repaired.append(_repair_slide(item, index, session))
    return repaired


def _solution_blocks(slides: list[ClassroomSlide]) -> list[tuple[ClassroomSlide, dict]]:
    out: list[tuple[ClassroomSlide, dict]] = []
    for item in slides:
        for block in _slide_blocks(item):
            btype = _visual_block_type(block)
            if btype in {"example", "exercise"} and (
                block.get("question") or block.get("answer") or block.get("explanation") or block.get("steps")
            ):
                out.append((item, block))
            if len(out) >= 4:
                return out
    return out


def _add_solution_slide(prs: Presentation, source: ClassroomSlide, block: dict, index: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    accent = _accent(source)
    header_item = source.model_copy(update={"kicker": f"解析 / {index:02d}"})
    _add_header(slide, header_item, index, total, accent)
    _textbox(slide, Inches(0.65), Inches(0.9), Inches(8.6), Inches(0.58), _clean(block.get("title") or source.title or "详细解析", 64), 25, INK, True)
    difficulty = _clean(block.get("difficulty") or block.get("level") or "解析", 32)
    _textbox(slide, Inches(10.05), Inches(0.98), Inches(2.25), Inches(0.26), difficulty, 10, accent, True, PP_ALIGN.RIGHT)

    question = _clean(block.get("question") or source.title, 170)
    _callout(slide, Inches(0.72), Inches(1.72), Inches(5.75), Inches(0.86), "题目", question, accent)

    formula = _latex_value(block)
    if formula:
        _formula_card(slide, Inches(6.85), Inches(1.72), Inches(5.55), Inches(2.1), block, source, accent, include_steps=False)
        steps_top = Inches(3.15)
        steps_h = Inches(3.0)
    else:
        steps_top = Inches(2.92)
        steps_h = Inches(3.25)

    steps = _split_items(block.get("steps") or block.get("items") or source.board, 6)
    _section_label(slide, Inches(0.72), steps_top - Inches(0.4), "解析步骤", accent)
    _steps(slide, Inches(0.72), steps_top, Inches(5.75), steps_h, steps, accent, columns=1)

    explanation = _clean(block.get("explanation") or block.get("analysis") or block.get("solution") or source.teacher_note, 190)
    answer = _clean(block.get("answer") or explanation, 140)
    _callout(slide, Inches(6.85), Inches(4.12), Inches(5.55), Inches(0.82), "关键解释", explanation or "重点是说明每一步使用的条件，而不是只写最终答案。", accent)
    _callout(slide, Inches(6.85), Inches(5.25), Inches(5.55), Inches(0.72), "答案检查", answer or "能独立复现推导，并解释结果含义。", accent)


def build_classroom_pptx(session: ClassroomSessionResponse) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    lesson_slides = _prepare_slides(session)
    solution_blocks = _solution_blocks(lesson_slides)
    total = len(lesson_slides) + len(solution_blocks) + 2
    _add_cover(prs, session, total)

    for index, item in enumerate(lesson_slides, start=2):
        kind = _slide_kind(item)
        if _content_score(item) < 4:
            item = _repair_slide(item, index, session)
        if kind == "example":
            _add_example_slide(prs, item, index, total)
        elif kind == "table":
            _add_table_slide(prs, item, index, total)
        elif kind == "quiz":
            _add_quiz_slide(prs, item, index, total)
        elif kind == "process":
            _add_process_slide(prs, item, index, total)
        else:
            _add_concept_slide(prs, item, index, total)

    solution_start = len(lesson_slides) + 2
    for offset, (source, block) in enumerate(solution_blocks):
        _add_solution_slide(prs, source, block, solution_start + offset, total)

    _add_summary_slide(prs, session, total)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()

"""从常见文件格式提取纯文本。"""



from __future__ import annotations



import io
import logging
import re
import unicodedata

from pathlib import Path



TEXT_EXTENSIONS = {

    ".md",

    ".markdown",

    ".txt",

    ".csv",

    ".json",

    ".py",

    ".java",

    ".c",

    ".cpp",

    ".h",

    ".hpp",

    ".html",

    ".htm",

    ".xml",

    ".yaml",

    ".yml",

    ".rst",

    ".tex",

    ".sql",

    ".js",

    ".ts",

    ".tsx",

    ".jsx",

    ".go",

    ".rs",

    ".php",

    ".rb",

    ".swift",

    ".kt",

    ".log",

    ".ini",

    ".cfg",

}



BINARY_EXTENSIONS = {

    ".pdf",

    ".docx",

    ".doc",

    ".pptx",

    ".ppt",

    ".xlsx",

    ".xls",

}

_SYMBOL_RUN_RE = re.compile(r"[!\"#$%&'()*+,./:;<=>?@\[\\\]^_`{|}~]{4,}")
_CONTROL_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")


def clean_extracted_text(text: str) -> str:
    """清理 PDF/PPT 抽取出的装饰符号、兼容字形和控制字符。"""
    if not text:
        return ""
    text = unicodedata.normalize("NFKC", text)
    text = _CONTROL_RE.sub(" ", text)
    text = text.replace("\ufeff", " ").replace("\u00a0", " ")
    text = text.replace("Ø", "\n- ")
    text = _SYMBOL_RUN_RE.sub(" ", text)
    text = re.sub(r"(^|[\s。；;，,])([upr])(?=[\u4e00-\u9fffA-Z])", r"\1- ", text)
    text = re.sub(r"(?<=[\u4e00-\u9fff)）:：])([upr])(?=[\u4e00-\u9fffA-Za-z])", " - ", text)
    text = re.sub(r"\s+\*\s+", " ", text)
    lines: list[str] = []
    for raw in text.splitlines():
        line = re.sub(r"\s+", " ", raw).strip()
        if not line:
            continue
        visible = re.sub(r"\s+", "", line)
        alpha_num = sum(1 for ch in visible if ch.isalnum() or "\u4e00" <= ch <= "\u9fff")
        if len(visible) >= 6 and alpha_num / max(len(visible), 1) < 0.35:
            continue
        lines.append(line)
    return "\n".join(lines).strip()





def supported_extensions() -> list[str]:

    return sorted(TEXT_EXTENSIONS | BINARY_EXTENSIONS)





def _extract_pptx(data: bytes) -> str:

    from pptx import Presentation



    prs = Presentation(io.BytesIO(data))

    parts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):

        slide_texts: list[str] = []

        for shape in slide.shapes:

            if not hasattr(shape, "text"):

                continue

            text = (shape.text or "").strip()

            if text:

                slide_texts.append(text)

        if slide_texts:

            parts.append(f"【幻灯片 {slide_idx}】\n" + "\n".join(slide_texts))

    return "\n\n".join(parts).strip()





def _extract_xlsx(data: bytes) -> str:

    from openpyxl import load_workbook



    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)

    parts: list[str] = []

    for sheet in wb.worksheets:

        rows: list[str] = []

        for row in sheet.iter_rows(max_row=500, values_only=True):

            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]

            if cells:

                rows.append(" | ".join(cells))

        if rows:

            parts.append(f"【工作表 {sheet.title}】\n" + "\n".join(rows[:200]))

    wb.close()

    return "\n\n".join(parts).strip()





def extract_text_from_bytes(filename: str, data: bytes) -> str:

    ext = Path(filename).suffix.lower()

    if ext in TEXT_EXTENSIONS:

        for enc in ("utf-8", "utf-8-sig", "gbk", "latin-1"):

            try:

                return clean_extracted_text(data.decode(enc))

            except UnicodeDecodeError:

                continue

        return clean_extracted_text(data.decode("utf-8", errors="replace"))



    if ext == ".pdf":

        try:
            logging.getLogger("pypdf").setLevel(logging.ERROR)

            from pypdf import PdfReader



            reader = PdfReader(io.BytesIO(data))

            parts = []

            for page in reader.pages:

                parts.append(page.extract_text() or "")

            return clean_extracted_text("\n\n".join(parts))

        except Exception as e:

            raise ValueError(f"PDF 解析失败: {e}") from e



    if ext in {".docx", ".doc"}:

        if ext == ".doc":

            raise ValueError("旧版 .doc 请用 Word 另存为 .docx 后上传")

        try:

            from docx import Document



            doc = Document(io.BytesIO(data))

            return clean_extracted_text("\n".join(p.text for p in doc.paragraphs if p.text.strip()))

        except Exception as e:

            raise ValueError(f"Word 解析失败: {e}") from e



    if ext in {".pptx", ".ppt"}:

        if ext == ".ppt":

            raise ValueError("旧版 .ppt 请用 PowerPoint 另存为 .pptx 后上传")

        try:

            text = _extract_pptx(data)

            if not text:

                raise ValueError("PPT 中未提取到文本（可能为纯图片幻灯片）")

            return clean_extracted_text(text)

        except ValueError:

            raise

        except Exception as e:

            raise ValueError(f"PPT 解析失败: {e}") from e



    if ext in {".xlsx", ".xls"}:

        if ext == ".xls":

            raise ValueError("旧版 .xls 请用 Excel 另存为 .xlsx 后上传")

        try:

            text = _extract_xlsx(data)

            if not text:

                raise ValueError("Excel 中未提取到有效单元格文本")

            return clean_extracted_text(text)

        except ValueError:

            raise

        except Exception as e:

            raise ValueError(f"Excel 解析失败: {e}") from e



    raise ValueError(f"不支持的文件类型: {ext or '(无扩展名)'}")





def guess_mime(filename: str) -> str:

    ext = Path(filename).suffix.lower()

    mapping = {

        ".pdf": "application/pdf",

        ".doc": "application/msword",

        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",

        ".ppt": "application/vnd.ms-powerpoint",

        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",

        ".xls": "application/vnd.ms-excel",

        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",

        ".md": "text/markdown",

        ".markdown": "text/markdown",

        ".txt": "text/plain",

        ".csv": "text/csv",

        ".html": "text/html",

        ".htm": "text/html",

        ".json": "application/json",

        ".py": "text/x-python",

    }

    return mapping.get(ext, "application/octet-stream")
